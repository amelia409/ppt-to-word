# 导入必要的库和模块
from pptx import Presentation  # 用于读取和处理PPT文件
import pytesseract  # OCR工具，用于从图像中提取文本
from PIL import Image  # 图像处理库
import io  # 用于处理二进制流
import re  # 正则表达式库，用于文本处理
import os  # 操作系统相关功能

def extract_structured_content(file_path: str) -> list:
    """从PPT文件中提取结构化内容
    
    参数:
        file_path: PPT文件的路径
        
    返回:
        包含所有幻灯片结构化内容的列表，每个幻灯片包含标题、内容和图片
    """
    prs = Presentation(file_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "slide_number": i+1,
            "title": "",
            "content": [],
            "images": []
        }
        
        for shape in slide.shapes:
            try:
                # 处理文本元素
                if hasattr(shape, "text") and shape.text.strip():  # 检查形状是否包含文本且不为空
                    text = clean_text(shape.text)  # 清理和规范化文本
                    if is_title(shape, slide):  # 判断是否为标题
                        slide_data["title"] = text  # 如果是标题，存储到标题字段
                    else:
                        # 如果不是标题，添加到内容列表中
                        slide_data["content"].append({
                            "type": "text",  # 内容类型为文本
                            "data": text  # 文本内容
                        })
                
                # 处理图片元素
                elif hasattr(shape, "shape_type") and shape.shape_type == 13:  # 13表示图片类型
                    try:
                        # 从形状中提取图片数据并转换为PIL图像对象
                        img = Image.open(io.BytesIO(shape.image.blob))
                        # 尝试从图像中提取文本（OCR识别）
                        img_text = extract_text_from_image(img)
                        # 将图片转换为二进制数据以便存储
                        img_byte_arr = io.BytesIO()
                        img.save(img_byte_arr, format='PNG')
                        # 将图片信息添加到幻灯片数据中
                        slide_data["images"].append({
                            "image": img_byte_arr.getvalue(),  # 图片二进制数据
                            "description": get_shape_description(shape),  # 图片描述
                            "text": img_text  # 从图片中提取的文本
                        })
                    except Exception as e:
                        print(f"处理图片时出错: {e}")  # 记录图片处理错误
                        
            except Exception as e:
                print(f"处理形状元素时出错: {e}")  # 记录形状处理错误
                continue  # 继续处理下一个形状元素
        
        slides_data.append(slide_data)
    return slides_data

def clean_text(text: str) -> str:
    """清理和规范化提取的文本"""
    # 处理中英文混合的换行问题
    text = re.sub(r'([\u4e00-\u9fa5a-zA-Z0-9])\s*\n\s*([\u4e00-\u9fa5a-zA-Z0-9])', r'\1 \2', text)
    # 规范化中文标点
    text = re.sub(r'([，。！？；：])\s+', r'\1', text)
    return text

def is_title(shape, slide) -> bool:
    """判断一个形状元素是否为标题
    
    使用多种方法来确定形状是否为标题，提高判断的准确性
    
    参数:
        shape: PPT中的形状元素
        slide: 形状所在的幻灯片
        
    返回:
        布尔值，表示该形状是否为标题
    """
    try:
        # 第一张幻灯片通常是标题幻灯片
        if slide.slide_index == 0:
            return True
            
        # 通过形状名称判断是否为标题
        if hasattr(shape, 'name') and 'title' in shape.name.lower():
            return True
            
        # 检查是否为占位符（这种方法与多个版本兼容）
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            return True
            
        return False  # 如果以上条件都不满足，则不是标题
    except:
        return False

def get_shape_description(shape) -> str:
    """获取形状的替代文本"""
    try:
        return shape.name if hasattr(shape, 'name') and shape.name else "图片"
    except:
        return "图片"
        
def extract_text_from_image(image):
    """从图像中提取文本，支持中文"""
    try:
        # 设置Tesseract OCR语言为中文简体+英文
        text = pytesseract.image_to_string(image, lang='chi_sim+eng')
        return text.strip() if text else ""
    except Exception as e:
        print(f"OCR识别错误: {e}")
        return ""